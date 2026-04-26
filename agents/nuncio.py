import anthropic
import argparse
import datetime
import os
import io
import json
import requests
import subprocess
import time
from bs4 import BeautifulSoup
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseUpload

CONFIRMATION_REQUIRED = {"send_email", "send_email_with_attachments", "create_calendar_event", "create_multiple_events", "upload_to_drive", "remember", "delete_memory"}

# --- Config ---
CREDS_FILE = os.path.join(os.path.dirname(__file__), '..', 'keys', 'google_credentials.json')
TOKEN_FILE = os.path.join(os.path.dirname(__file__), '..', 'keys', 'google_token.json')
TELEGRAM_KEY_FILE = os.path.join(os.path.dirname(__file__), '..', 'keys', 'telegram.key')
TELEGRAM_USER_ID_FILE = os.path.join(os.path.dirname(__file__), '..', 'keys', 'telegram_user_id.key')

SCOPES = [
    'https://www.googleapis.com/auth/calendar',
    'https://www.googleapis.com/auth/gmail.readonly',
    'https://www.googleapis.com/auth/gmail.send',
    'https://www.googleapis.com/auth/drive.file',
    'https://www.googleapis.com/auth/drive.readonly'
]

NUNCIO_FOLDER = os.path.join(os.path.dirname(__file__), '..', 'nuncio-inbox')

# --- Error classification helper ---
def _classify_error(e):
    """Returns (reason, retryable) for structured error responses."""
    if isinstance(e, (FileNotFoundError, IsADirectoryError)):
        return "file_not_found", False
    if isinstance(e, PermissionError):
        return "permission_error", False
    if isinstance(e, (requests.exceptions.Timeout,)):
        return "network_or_timeout_error", True
    if isinstance(e, requests.exceptions.ConnectionError):
        return "network_or_timeout_error", True
    if isinstance(e, requests.exceptions.HTTPError):
        status = e.response.status_code if e.response is not None else 0
        if 400 <= status < 500:
            return f"http_{status}_error", False
        return f"http_{status}_error", True
    try:
        from googleapiclient.errors import HttpError
        if isinstance(e, HttpError):
            status = int(e.resp.status)
            if 400 <= status < 500:
                return f"http_{status}_error", False
            return f"http_{status}_error", True
    except ImportError:
        pass
    return "unknown_error", True

# --- History ---

HISTORY_FILE = os.path.join(os.path.dirname(__file__), '..', 'logs', 'conversation_history.json')
BOOK_SCOUT_FILE = os.path.join(os.path.dirname(__file__), '..', 'logs', 'book_scout_last_run.txt')
BOOK_PREFERENCES_FILE = os.path.join(os.path.dirname(__file__), 'book_preferences.md')

def load_history():
    if os.path.exists(HISTORY_FILE):
        try:
            with open(HISTORY_FILE, 'r') as f:
                history = json.load(f)
            for msg in history:
                if not isinstance(msg, dict) or 'role' not in msg or 'content' not in msg:
                    raise ValueError("Malformed message in history")
                if isinstance(msg['content'], str):
                    raise ValueError("Content is a plain string, not a list of blocks")
            return history
        except (json.JSONDecodeError, ValueError) as e:
            print(f"[Nuncio] Warning: conversation history was corrupted and has been reset. ({e})")
            save_history([])
            return []
    return []

def save_history(history):
    with open(HISTORY_FILE, 'w') as f:
        json.dump(history, f, indent=2, default=str)


# --- Book Scout ---

def get_book_scout_status():
    if not os.path.exists(BOOK_SCOUT_FILE):
        return "never"
    with open(BOOK_SCOUT_FILE, 'r') as f:
        return f.read().strip()

def save_book_scout_timestamp():
    with open(BOOK_SCOUT_FILE, 'w') as f:
        f.write(datetime.datetime.now().strftime("%Y-%m-%d"))

def book_scout_prompt_fragment():
    last_run = get_book_scout_status()
    if last_run == "never":
        return "The book scout has never been run."
    try:
        last_date = datetime.datetime.strptime(last_run, "%Y-%m-%d")
        days_ago = (datetime.datetime.now() - last_date).days
        if days_ago >= 7:
            return f"The book scout was last run {days_ago} days ago (on {last_run}). This is more than a week — proactively ask Vernie at the start of this conversation if she would like you to run it now."
        else:
            return f"The book scout was last run {days_ago} days ago (on {last_run}). No need to prompt yet."
    except ValueError:
        return "The book scout has never been run."



# --- Google Auth ---
def get_credentials():
    from google.auth.transport.requests import Request
    creds = None
    if os.path.exists(TOKEN_FILE):
        creds = Credentials.from_authorized_user_file(TOKEN_FILE, SCOPES)
    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            creds.refresh(Request())
            with open(TOKEN_FILE, 'w') as token:
                token.write(creds.to_json())
        else:
            flow = InstalledAppFlow.from_client_secrets_file(CREDS_FILE, SCOPES)
            creds = flow.run_local_server(port=0)
            with open(TOKEN_FILE, 'w') as token:
                token.write(creds.to_json())
    return creds

# --- Calendar Tool ---
def get_calendar_events(query=None):
    try:
        creds = get_credentials()
        service = build('calendar', 'v3', credentials=creds)
        now = datetime.datetime.utcnow().isoformat() + 'Z'
        events_result = service.events().list(
            calendarId='primary',
            timeMin=now,
            maxResults=20,
            singleEvents=True,
            orderBy='startTime'
        ).execute()
        events = events_result.get('items', [])
        if query:
            events = [e for e in events if query.lower() in e.get('summary', '').lower() or query.lower() in e.get('description', '').lower()]
        if not events:
            return "No upcoming events found."
        result = ""
        for event in events:
            start = event['start'].get('dateTime', event['start'].get('date'))
            result += f"{start}: {event['summary']}\n"
        return result
    except Exception as e:
        reason, retryable = _classify_error(e)
        return json.dumps({"status": "error", "reason": reason, "retryable": retryable, "detail": str(e)})

def create_calendar_event(summary, start_datetime, end_datetime, description=None):
    creds = get_credentials()
    service = build('calendar', 'v3', credentials=creds)
    event = {
        'summary': f'[Nuncio] {summary}',
        'start': {'dateTime': start_datetime, 'timeZone': 'Asia/Singapore'},
        'end': {'dateTime': end_datetime, 'timeZone': 'Asia/Singapore'},
    }
    if description:
        event['description'] = description
    created = service.events().insert(calendarId='primary', body=event).execute()
    return f"Event created: {created['summary']} on {start_datetime}"

def create_multiple_events(events_list):
    results = []
    for event in events_list:
        try:
            result = create_calendar_event(
                event['summary'],
                event['start_datetime'],
                event['end_datetime'],
                event.get('description')
            )
            results.append(f"✓ {result}")
        except Exception as e:
            results.append(f"✗ Failed to create '{event['summary']}': {str(e)}")
    return "\n".join(results)

# --- Gmail Tool ---
def get_recent_emails(query=None):
    try:
        creds = get_credentials()
        service = build('gmail', 'v1', credentials=creds)
        q = query if query else ''
        results = service.users().messages().list(userId='me', maxResults=5, q=q).execute()
        messages = results.get('messages', [])
        if not messages:
            return "No emails found."
        result = ""
        for msg in messages:
            txt = service.users().messages().get(userId='me', id=msg['id']).execute()
            headers = txt['payload']['headers']
            subject = next((h['value'] for h in headers if h['name'] == 'Subject'), 'No Subject')
            sender = next((h['value'] for h in headers if h['name'] == 'From'), 'Unknown')
            result += f"From: {sender}\nSubject: {subject}\n\n"
        return result
    except Exception as e:
        reason, retryable = _classify_error(e)
        return json.dumps({"status": "error", "reason": reason, "retryable": retryable, "detail": str(e)})

def send_telegram_message(text):
    try:
        with open(TELEGRAM_KEY_FILE) as f:
            token = f.read().strip()
        with open(TELEGRAM_USER_ID_FILE) as f:
            user_id = f.read().strip()
        url = f"https://api.telegram.org/bot{token}/sendMessage"
        resp = requests.post(url, json={"chat_id": user_id, "text": text}, timeout=10)
        resp.raise_for_status()
        return "Telegram message sent successfully."
    except FileNotFoundError:
        return "Error: Telegram not configured (key files missing)."
    except Exception as e:
        reason, retryable = _classify_error(e)
        return json.dumps({"status": "error", "reason": reason, "retryable": retryable, "detail": str(e)})

def send_email(to, subject, body):
    import base64
    from email.mime.text import MIMEText
    creds = get_credentials()
    service = build('gmail', 'v1', credentials=creds)
    message = MIMEText(body)
    message['to'] = to
    message['subject'] = subject
    raw = base64.urlsafe_b64encode(message.as_bytes()).decode()
    service.users().messages().send(userId='me', body={'raw': raw}).execute()
    return f"Email sent to {to} with subject '{subject}'"

def send_email_with_attachments(to, subject, body, attachment_paths):
    import base64
    import mimetypes
    from email.mime.multipart import MIMEMultipart
    from email.mime.text import MIMEText
    from email.mime.base import MIMEBase
    from email import encoders
    try:
        creds = get_credentials()
        service = build('gmail', 'v1', credentials=creds)
        message = MIMEMultipart()
        message['to'] = to
        message['subject'] = subject
        message.attach(MIMEText(body))
        missing = []
        for filepath in attachment_paths:
            if not os.path.exists(filepath):
                missing.append(filepath)
                continue
            mimetype, _ = mimetypes.guess_type(filepath)
            if mimetype is None:
                _, ext = os.path.splitext(filepath)
                mimetype = MIMETYPE_MAP.get(ext.lower(), 'application/octet-stream')
            maintype, subtype = mimetype.split('/', 1)
            with open(filepath, 'rb') as fh:
                part = MIMEBase(maintype, subtype)
                part.set_payload(fh.read())
            encoders.encode_base64(part)
            part.add_header('Content-Disposition', 'attachment', filename=os.path.basename(filepath))
            message.attach(part)
        if missing:
            return json.dumps({"status": "error", "reason": "file_not_found", "retryable": False, "detail": f"Files not found: {missing}"})
        raw = base64.urlsafe_b64encode(message.as_bytes()).decode()
        service.users().messages().send(userId='me', body={'raw': raw}).execute()
        names = [os.path.basename(p) for p in attachment_paths]
        return f"Email sent to {to} with subject '{subject}' and attachments: {', '.join(names)}"
    except Exception as e:
        reason, retryable = _classify_error(e)
        return json.dumps({"status": "error", "reason": reason, "retryable": retryable, "detail": str(e)})

# --- Drive Tools ---
def list_drive_files(query=None):
    try:
        creds = get_credentials()
        service = build('drive', 'v3', credentials=creds)
        q = f"name contains '{query}'" if query else ""
        results = service.files().list(
            pageSize=20,
            orderBy='modifiedTime desc',
            q=q,
            fields="files(id, name, mimeType, modifiedTime)"
        ).execute()
        files = results.get('files', [])
        if not files:
            return "No files found."
        result = f"Found {len(files)} files:\n\n"
        for f in files:
            result += f"{f['modifiedTime']}: {f['name']} ({f['mimeType']})\n"
        return result
    except Exception as e:
        reason, retryable = _classify_error(e)
        return json.dumps({"status": "error", "reason": reason, "retryable": retryable, "detail": str(e)})

def upload_to_drive(filename, content):
    creds = get_credentials()
    service = build('drive', 'v3', credentials=creds)
    file_metadata = {'name': filename}
    media = MediaIoBaseUpload(
        io.BytesIO(content.encode('utf-8')),
        mimetype='text/plain'
    )
    file = service.files().create(
        body=file_metadata,
        media_body=media,
        fields='id, name'
    ).execute()
    return f"File '{file['name']}' uploaded to Google Drive with ID: {file['id']}"

MIMETYPE_MAP = {
    '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
    '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
    '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
    '.pdf':  'application/pdf',
    '.txt':  'text/plain',
    '.csv':  'text/csv',
    '.png':  'image/png',
    '.jpg':  'image/jpeg',
    '.jpeg': 'image/jpeg',
}

def upload_file_to_drive(filepath):
    import mimetypes
    try:
        if not os.path.exists(filepath):
            return json.dumps({"status": "error", "reason": "file_not_found", "retryable": False, "detail": f"File not found: {filepath}"})
        _, ext = os.path.splitext(filepath)
        mimetype = MIMETYPE_MAP.get(ext.lower()) or mimetypes.guess_type(filepath)[0] or 'application/octet-stream'
        filename = os.path.basename(filepath)
        creds = get_credentials()
        service = build('drive', 'v3', credentials=creds)
        file_metadata = {'name': filename}
        with open(filepath, 'rb') as fh:
            media = MediaIoBaseUpload(fh, mimetype=mimetype, resumable=True)
            file = service.files().create(
                body=file_metadata,
                media_body=media,
                fields='id, name'
            ).execute()
        return f"File '{file['name']}' uploaded to Google Drive with ID: {file['id']}"
    except Exception as e:
        reason, retryable = _classify_error(e)
        return json.dumps({"status": "error", "reason": reason, "retryable": retryable, "detail": str(e)})


# --- Local File Tools ---

def list_files():
    result = []
    for root, dirs, files in os.walk(NUNCIO_FOLDER):
        for file in files:
            full = os.path.join(root, file)
            relative = os.path.relpath(full, NUNCIO_FOLDER)
            result.append(relative)
    return "\n".join(result) if result else "No files found."

def write_file(filepath, content):
    with open(filepath, 'w') as f:
        f.write(content)
    return f"File written to {filepath}"

def read_file(filepath):
    try:
        if os.path.isdir(filepath):
            return json.dumps({"status": "error", "reason": "path_is_directory", "retryable": False, "detail": f"'{filepath}' is a directory, not a file. Use list_files to see its contents."})
        with open(filepath, 'r') as f:
            return f.read()
    except FileNotFoundError as e:
        return json.dumps({"status": "error", "reason": "file_not_found", "retryable": False, "detail": str(e)})
    except PermissionError as e:
        return json.dumps({"status": "error", "reason": "permission_error", "retryable": False, "detail": str(e)})
    except Exception as e:
        reason, retryable = _classify_error(e)
        return json.dumps({"status": "error", "reason": reason, "retryable": retryable, "detail": str(e)})
    
def read_pdf(filepath):
    import fitz
    try:
        if not os.path.exists(filepath):
            return json.dumps({"status": "error", "reason": "file_not_found", "retryable": False, "detail": f"File not found: {filepath}"})
        doc = fitz.open(filepath)
        result = ""
        for page in doc:
            result += page.get_text()
        return result
    except PermissionError as e:
        return json.dumps({"status": "error", "reason": "permission_error", "retryable": False, "detail": str(e)})
    except Exception as e:
        reason, retryable = _classify_error(e)
        return json.dumps({"status": "error", "reason": reason, "retryable": retryable, "detail": str(e)})

def read_docx(filepath):
    from docx import Document
    try:
        if not os.path.exists(filepath):
            return json.dumps({"status": "error", "reason": "file_not_found", "retryable": False, "detail": f"File not found: {filepath}"})
        doc = Document(filepath)
        return "\n".join([para.text for para in doc.paragraphs])
    except PermissionError as e:
        return json.dumps({"status": "error", "reason": "permission_error", "retryable": False, "detail": str(e)})
    except Exception as e:
        reason, retryable = _classify_error(e)
        return json.dumps({"status": "error", "reason": reason, "retryable": retryable, "detail": str(e)})

# --- Presentation Tool ---

def create_presentation(title, slides, output_filename=None):
    from pptx import Presentation
    from pptx.util import Pt
    try:
        prs = Presentation()

        # Title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = title
        if len(title_slide.placeholders) > 1:
            title_slide.placeholders[1].text = ""

        # Content slides
        for s in slides:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = s.get("title", "")
            tf = slide.placeholders[1].text_frame
            tf.clear()
            bullets = s.get("bullets", [])
            for i, bullet in enumerate(bullets):
                if i == 0:
                    tf.paragraphs[0].text = bullet
                else:
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.level = s.get("level", 0) if isinstance(bullet, str) else bullet.get("level", 0)
            if s.get("notes"):
                slide.notes_slide.notes_text_frame.text = s["notes"]

        if output_filename is None:
            slug = title.lower().replace(" ", "_")
            slug = "".join(c for c in slug if c.isalnum() or c == "_")
            output_filename = slug + ".pptx"

        output_path = os.path.join(NUNCIO_FOLDER, output_filename)
        prs.save(output_path)
        return f"Presentation saved to {output_path} ({len(slides)} slides). Use upload_file_to_drive with filepath='{output_path}' to share it on Google Drive."
    except Exception as e:
        reason, retryable = _classify_error(e)
        return json.dumps({"status": "error", "reason": reason, "retryable": retryable, "detail": str(e)})


# --- Browser session (persistent across tool calls) ---
_playwright_instance = None
_browser_instance = None
_page_instance = None
_browser_owned = False  # True if Nuncio launched the browser; False if attached via CDP

# Chrome default: 9222  |  Firefox default: 9223
CDP_PORTS = [9222, 9223]

_cdp_attach_failed_reason = None  # set when CDP attach is attempted but fails

# --- Firefox CDP session ---
_firefox_cdp_session = None


class FirefoxCDPSession:
    """CDP-over-WebSocket bridge to an existing Firefox with --remote-debugging-port."""

    def __init__(self, port, ws_url):
        import websocket as _ws
        self.port = port
        self._ws = _ws.WebSocket()
        self._ws.settimeout(15)
        self._ws.connect(ws_url)
        self._msg_id = 0

    def _send(self, method, params=None):
        self._msg_id += 1
        mid = self._msg_id
        self._ws.send(json.dumps({"id": mid, "method": method, "params": params or {}}))
        deadline = time.time() + 15
        while time.time() < deadline:
            try:
                self._ws.settimeout(max(0.1, deadline - time.time()))
                raw = self._ws.recv()
            except Exception:
                break
            if not raw:
                continue
            data = json.loads(raw)
            if data.get("id") == mid:
                if "error" in data:
                    raise RuntimeError(f"CDP: {data['error'].get('message', str(data['error']))}")
                return data.get("result", {})
        raise TimeoutError(f"No CDP response for {method}")

    def evaluate(self, expression):
        result = self._send("Runtime.evaluate", {
            "expression": expression,
            "returnByValue": True,
            "awaitPromise": False,
        })
        val = result.get("result", {})
        if val.get("type") == "undefined" or val.get("subtype") == "null":
            return None
        return val.get("value")

    def navigate(self, url):
        self._send("Page.navigate", {"url": url})
        self._wait_for_load()

    def _wait_for_load(self, timeout=15):
        start = time.time()
        while time.time() - start < timeout:
            try:
                if self.evaluate("document.readyState") == "complete":
                    return
            except Exception:
                pass
            time.sleep(0.4)

    def title(self):
        return self.evaluate("document.title") or ""

    def current_url(self):
        return self.evaluate("window.location.href") or ""

    def reconnect(self, ws_url):
        try:
            self._ws.close()
        except Exception:
            pass
        import websocket as _ws
        self._ws = _ws.WebSocket()
        self._ws.settimeout(15)
        self._ws.connect(ws_url)

    def close(self):
        try:
            self._ws.close()
        except Exception:
            pass


def _detect_browser_on_port(port):
    """Returns 'firefox', 'chrome', or None based on /json/version."""
    try:
        resp = requests.get(f"http://localhost:{port}/json/version", timeout=1)
        browser = resp.json().get("Browser", "")
        if "Firefox" in browser:
            return "firefox"
        if "Chrome" in browser or "Chromium" in browser:
            return "chrome"
    except Exception:
        pass
    return None


def _firefox_list_tabs(port):
    """Return page/tab target dicts from Firefox CDP."""
    try:
        resp = requests.get(f"http://localhost:{port}/json/list", timeout=2)
        return [t for t in resp.json()
                if t.get("webSocketDebuggerUrl") and t.get("type") in ("tab", "page")]
    except Exception:
        return []


def _get_firefox_session():
    """Return active Firefox CDP session, auto-connecting if Firefox is on a known debug port."""
    global _firefox_cdp_session
    if _firefox_cdp_session is not None:
        try:
            _firefox_cdp_session.evaluate("1")
            return _firefox_cdp_session
        except Exception:
            _firefox_cdp_session.close()
            _firefox_cdp_session = None
    import socket as _sock
    for port in CDP_PORTS:
        try:
            s = _sock.create_connection(("127.0.0.1", port), timeout=0.5)
            s.close()
        except OSError:
            continue
        if _detect_browser_on_port(port) == "firefox":
            tabs = _firefox_list_tabs(port)
            if tabs:
                _firefox_cdp_session = FirefoxCDPSession(port, tabs[0]["webSocketDebuggerUrl"])
                return _firefox_cdp_session
    return None


def _get_page():
    global _playwright_instance, _browser_instance, _page_instance, _browser_owned, _cdp_attach_failed_reason
    import socket
    from playwright.sync_api import sync_playwright
    if _playwright_instance is None:
        _playwright_instance = sync_playwright().start()
        connected = False
        for port in CDP_PORTS:
            # Quick TCP probe — avoids a slow Playwright timeout on a closed port
            try:
                s = socket.create_connection(("127.0.0.1", port), timeout=0.5)
                s.close()
            except OSError:
                continue
            # Skip Firefox ports — handled via FirefoxCDPSession, not Playwright
            if _detect_browser_on_port(port) == "firefox":
                continue
            try:
                _browser_instance = _playwright_instance.chromium.connect_over_cdp(f"http://localhost:{port}")
                _browser_owned = False
                contexts = _browser_instance.contexts
                if contexts and contexts[0].pages:
                    _page_instance = contexts[0].pages[0]
                else:
                    ctx = contexts[0] if contexts else _browser_instance.new_context()
                    _page_instance = ctx.new_page()
                connected = True
                _cdp_attach_failed_reason = None
                break
            except Exception as e:
                _cdp_attach_failed_reason = f"port {port} is open but CDP attach failed: {e}"
        if not connected:
            if not _cdp_attach_failed_reason:
                _cdp_attach_failed_reason = f"No Chrome/Chromium found on ports {CDP_PORTS}. Start Chrome with: google-chrome --remote-debugging-port=9222"
            _browser_instance = _playwright_instance.chromium.launch(headless=False)
            _page_instance = _browser_instance.new_page()
            _browser_owned = True
    return _page_instance

def _match_domain(page_url, target_url):
    from urllib.parse import urlparse
    try:
        p = urlparse(page_url).netloc.lstrip("www.")
        t = urlparse(target_url).netloc.lstrip("www.")
        return p and t and (p == t or p.endswith("." + t) or t.endswith("." + p))
    except Exception:
        return False

def _get_domain(url):
    from urllib.parse import urlparse
    try:
        return urlparse(url).netloc.lstrip("www.").split(':')[0]
    except Exception:
        return ""

def _wmctrl_windows():
    """Return [(wid, title)] for all X11 windows via wmctrl, or [] if unavailable."""
    try:
        r = subprocess.run(['wmctrl', '-l'], capture_output=True, text=True, timeout=3)
        if r.returncode == 0:
            rows = []
            for line in r.stdout.splitlines():
                parts = line.split(None, 3)
                if len(parts) >= 4:
                    rows.append((parts[0], parts[3]))
            return rows
    except (FileNotFoundError, subprocess.TimeoutExpired, OSError):
        pass
    return []

def _browser_type_from_title(title):
    t = title.lower()
    if 'firefox' in t or 'mozilla firefox' in t:
        return 'firefox'
    if 'google chrome' in t or 'chromium' in t:
        return 'chrome'
    return None

def _find_domain_in_windows(domain, windows):
    """
    Search open windows for a browser whose title contains the domain.
    Returns (browser_type, wid, title) prioritising Firefox, or (None, None, None).
    """
    firefox_match = chrome_match = None
    for wid, title in windows:
        if domain.lower() not in title.lower():
            continue
        bt = _browser_type_from_title(title)
        if bt == 'firefox' and not firefox_match:
            firefox_match = ('firefox', wid, title)
        elif bt == 'chrome' and not chrome_match:
            chrome_match = ('chrome', wid, title)
    return firefox_match or chrome_match or (None, None, None)

def _first_running_browser(windows):
    """
    Return (browser_type, wid) for the first Firefox or Chrome window found
    (Firefox preferred). Falls back to pgrep if wmctrl found nothing.
    Returns (None, None) if no browser is running.
    """
    firefox = chrome = None
    for wid, title in windows:
        bt = _browser_type_from_title(title)
        if bt == 'firefox' and not firefox:
            firefox = ('firefox', wid)
        elif bt == 'chrome' and not chrome:
            chrome = ('chrome', wid)
    result = firefox or chrome
    if result:
        return result
    for pname, bt in [('firefox', 'firefox'), ('google-chrome', 'chrome'), ('chromium-browser', 'chrome')]:
        try:
            r = subprocess.run(['pgrep', '-f', pname], capture_output=True, timeout=2)
            if r.returncode == 0:
                return (bt, None)
        except Exception:
            pass
    return (None, None)

def _activate_window(wid):
    """Raise a window to the front via wmctrl or xdotool."""
    for cmd in [['wmctrl', '-ia', wid], ['xdotool', 'windowactivate', '--sync', wid]]:
        try:
            subprocess.run(cmd, timeout=3, check=True, capture_output=True)
            return True
        except Exception:
            continue
    return False

def browser_navigate(url):
    global _page_instance, _firefox_cdp_session, _browser_instance, _playwright_instance, _browser_owned, _cdp_attach_failed_reason
    try:
        import socket
        windows = _wmctrl_windows()
        domain = _get_domain(url)

        # 1. Firefox CDP (preferred) — attach to existing Firefox with --remote-debugging-port
        for port in CDP_PORTS:
            try:
                s = socket.create_connection(("127.0.0.1", port), timeout=0.5)
                s.close()
            except OSError:
                continue
            if _detect_browser_on_port(port) != "firefox":
                continue
            tabs = _firefox_list_tabs(port)
            target_tab = None
            if domain:
                for tab in tabs:
                    if domain in tab.get("url", ""):
                        target_tab = tab
                        break
            if target_tab:
                # Domain already open — connect to that tab
                if _firefox_cdp_session is not None:
                    _firefox_cdp_session.reconnect(target_tab["webSocketDebuggerUrl"])
                else:
                    _firefox_cdp_session = FirefoxCDPSession(port, target_tab["webSocketDebuggerUrl"])
                _, wid, _ = _find_domain_in_windows(domain, windows)
                if not wid:
                    for wid2, t in windows:
                        if "firefox" in t.lower():
                            wid = wid2
                            break
                if wid:
                    _activate_window(wid)
                return f"Switched to existing Firefox tab: {target_tab.get('title', '')} ({target_tab.get('url', '')})"
            # Domain not open — navigate in Firefox
            if _firefox_cdp_session is None:
                if not tabs:
                    return "Firefox is running with remote debugging but has no open tabs."
                _firefox_cdp_session = FirefoxCDPSession(port, tabs[0]["webSocketDebuggerUrl"])
            _firefox_cdp_session.navigate(url)
            title = _firefox_cdp_session.title()
            _, running_wid = _first_running_browser(windows)
            if running_wid:
                _activate_window(running_wid)
            return f"Navigated Firefox to: {title}"

        # 2. Playwright session already live (Chrome/Chromium)
        if _browser_instance is not None:
            if not _browser_owned:
                for ctx in _browser_instance.contexts:
                    for p in ctx.pages:
                        try:
                            if _match_domain(p.url, url):
                                _page_instance = p
                                p.bring_to_front()
                                return f"Switched to existing tab: {p.title()} ({p.url})"
                        except Exception:
                            pass
            _page_instance.goto(url, timeout=15000)
            return f"Navigated to: {_page_instance.title()}"

        # 3. Firefox running but without remote debugging
        running_bt, running_wid = _first_running_browser(windows)
        if running_bt == "firefox":
            try:
                subprocess.Popen(["firefox", url])
                if running_wid:
                    _activate_window(running_wid)
            except Exception:
                pass
            return (
                f"Opened {url} in Firefox, but interactive tools (fill, click, read) are unavailable. "
                "To enable full interaction, start Firefox with: firefox --remote-debugging-port=9222"
            )

        # 4. Chrome CDP or headless Chromium fallback
        _get_page()
        if not _browser_owned:
            for ctx in _browser_instance.contexts:
                for p in ctx.pages:
                    try:
                        if _match_domain(p.url, url):
                            _page_instance = p
                            p.bring_to_front()
                            return f"Switched to existing tab: {p.title()} ({p.url})"
                    except Exception:
                        pass
        _page_instance.goto(url, timeout=15000)
        note = ""
        if _browser_owned and _cdp_attach_failed_reason:
            note = f" [Note: {_cdp_attach_failed_reason}]"
        return f"Navigated to: {_page_instance.title()}{note}"

    except Exception as e:
        return f"Error: {str(e)}"

def browser_fill(label, text, selector=None):
    session = _get_firefox_session()
    if session is not None:
        lbl = json.dumps(label)
        val = json.dumps(text)
        if selector:
            sel = json.dumps(selector)
            expr = f"""(function(){{
                var el = document.querySelector({sel});
                if (!el) return 'not_found';
                el.focus(); el.value = {val};
                el.dispatchEvent(new Event('input', {{bubbles:true}}));
                el.dispatchEvent(new Event('change', {{bubbles:true}}));
                return 'filled';
            }})()"""
        else:
            expr = f"""(function(){{
                var label = {lbl}, val = {val}, found = null;
                for (var inp of document.querySelectorAll('input,textarea')) {{
                    if (inp.placeholder && inp.placeholder.toLowerCase() === label.toLowerCase()) {{ found = inp; break; }}
                }}
                if (!found) {{
                    for (var inp of document.querySelectorAll('[aria-label]')) {{
                        if (inp.getAttribute('aria-label').toLowerCase() === label.toLowerCase()) {{ found = inp; break; }}
                    }}
                }}
                if (!found) found = document.querySelector('input[name=' + JSON.stringify(label) + '],textarea[name=' + JSON.stringify(label) + ']');
                if (!found) {{
                    var el = document.getElementById(label);
                    if (el && (el.tagName === 'INPUT' || el.tagName === 'TEXTAREA')) found = el;
                }}
                if (!found) {{
                    for (var lbl of document.querySelectorAll('label')) {{
                        if (lbl.textContent.trim().toLowerCase().includes(label.toLowerCase())) {{
                            found = lbl.control || document.getElementById(lbl.htmlFor);
                            if (found) break;
                        }}
                    }}
                }}
                if (!found) found = document.querySelector('input[type="search"],input[type="text"],textarea');
                if (!found) return 'not_found';
                found.focus(); found.value = val;
                found.dispatchEvent(new Event('input', {{bubbles:true}}));
                found.dispatchEvent(new Event('change', {{bubbles:true}}));
                return 'filled';
            }})()"""
        result = session.evaluate(expr)
        if result == "filled":
            return f"Filled field with '{text}'"
        return f"Could not find field '{label}'. Use browser_inspect_forms to see available fields, then retry with selector=."
    page = _get_page()
    if selector:
        try:
            page.locator(selector).fill(text)
            return f"Filled '{selector}' with '{text}'"
        except Exception as e:
            return f"Could not fill '{selector}': {str(e)}"
    errors = []
    for attempt in [
        lambda: page.get_by_placeholder(label).fill(text),
        lambda: page.get_by_label(label).fill(text),
        lambda: page.locator(f"input[aria-label='{label}']").fill(text),
        lambda: page.locator(f"input[name='{label}']").fill(text),
        lambda: page.locator(f"input[id='{label}']").fill(text),
        lambda: page.locator(f"input[type='search']").first.fill(text),
        lambda: page.locator("input[type='text']").first.fill(text),
    ]:
        try:
            attempt()
            return f"Filled field with '{text}'"
        except Exception as e:
            errors.append(str(e))
    return f"Could not find field '{label}'. Use browser_inspect_forms to see available fields, then retry with selector=. Last error: {errors[-1]}"

def browser_click(text, selector=None):
    session = _get_firefox_session()
    if session is not None:
        txt = json.dumps(text)
        if selector:
            sel = json.dumps(selector)
            expr = f"""(function(){{
                var el = document.querySelector({sel});
                if (!el) return 'not_found';
                el.click();
                return 'clicked';
            }})()"""
        else:
            expr = f"""(function(){{
                var text = {txt};
                for (var el of document.querySelectorAll('button,input[type="submit"],input[type="button"],a,[role="button"]')) {{
                    var t = (el.textContent || el.value || '').trim();
                    if (t.toLowerCase().includes(text.toLowerCase())) {{ el.click(); return 'clicked'; }}
                }}
                return 'not_found';
            }})()"""
        result = session.evaluate(expr)
        if result == "clicked":
            return f"Clicked '{text}'"
        return f"Could not click '{text}'. Use browser_inspect_forms to see available buttons, then retry with selector=."
    page = _get_page()
    if selector:
        try:
            page.locator(selector).click()
            return f"Clicked '{selector}'"
        except Exception as e:
            return f"Could not click '{selector}': {str(e)}"
    errors = []
    for attempt in [
        lambda: page.get_by_role("button", name=text).click(),
        lambda: page.locator(f"input[type='submit'][value='{text}']").click(),
        lambda: page.locator("input[type='submit']").first.click(),
        lambda: page.get_by_text(text).first.click(),
    ]:
        try:
            attempt()
            return f"Clicked '{text}'"
        except Exception as e:
            errors.append(str(e))
    return f"Could not click '{text}'. Use browser_inspect_forms to see available buttons, then retry with selector=. Last error: {errors[-1]}"

def browser_inspect_forms():
    session = _get_firefox_session()
    if session is not None:
        try:
            data = session.evaluate("""Array.from(document.querySelectorAll('input,textarea,select,button,[role="button"]')).map(el => {
                var attrs = {};
                for (var a of el.attributes) attrs[a.name] = a.value;
                return {tag: el.tagName.toLowerCase(), attrs: attrs, text: (el.innerText||'').trim().slice(0,80)};
            })""")
            if not data:
                return "No form elements found on this page."
            lines = []
            for el in data:
                attrs_str = " ".join(f'{k}="{v}"' for k, v in el.get("attrs", {}).items())
                inner = f' (text: "{el["text"]}")' if el.get("text") else ""
                lines.append(f'<{el["tag"]} {attrs_str}>{inner}')
            return "\n".join(lines)
        except Exception as e:
            return f"Error inspecting forms: {str(e)}"
    try:
        page = _get_page()
        html = page.evaluate("""() => {
            const els = document.querySelectorAll('input, textarea, select, button, [role="button"]');
            return Array.from(els).map(el => {
                const attrs = {};
                for (const a of el.attributes) attrs[a.name] = a.value;
                return { tag: el.tagName.toLowerCase(), attrs, text: el.innerText?.trim().slice(0, 80) || '' };
            });
        }""")
        if not html:
            return "No form elements found on this page."
        lines = []
        for el in html:
            attrs_str = " ".join(f'{k}="{v}"' for k, v in el["attrs"].items())
            inner = f' (text: "{el["text"]}")' if el["text"] else ""
            lines.append(f'<{el["tag"]} {attrs_str}>{inner}')
        return "\n".join(lines)
    except Exception as e:
        return f"Error inspecting forms: {str(e)}"

def browser_read():
    session = _get_firefox_session()
    if session is not None:
        try:
            title = session.title()
            text = session.evaluate("document.body.innerText") or ""
            text = text[:5000] if len(text) > 5000 else text
            return f"Title: {title}\n\n{text}"
        except Exception as e:
            return f"Error reading page: {str(e)}"
    try:
        page = _get_page()
        title = page.title()
        text = page.inner_text("body")
        text = text[:5000] if len(text) > 5000 else text
        return f"Title: {title}\n\n{text}"
    except Exception as e:
        return f"Error reading page: {str(e)}"

def browser_close():
    global _playwright_instance, _browser_instance, _page_instance, _browser_owned, _firefox_cdp_session
    msgs = []
    if _firefox_cdp_session is not None:
        try:
            _firefox_cdp_session.close()
        except Exception:
            pass
        _firefox_cdp_session = None
        msgs.append("Disconnected from Firefox (your browser window is still open).")
    had_playwright = _playwright_instance is not None
    owned = _browser_owned
    try:
        if _browser_owned and _browser_instance:
            _browser_instance.close()
        if _playwright_instance:
            _playwright_instance.stop()
    except Exception as e:
        msgs.append(f"Error closing Chromium: {str(e)}")
    _playwright_instance = None
    _browser_instance = None
    _page_instance = None
    _browser_owned = False
    if had_playwright:
        msgs.append("Browser closed." if owned else "Disconnected from Chrome (your browser window is still open).")
    return " ".join(msgs) if msgs else "No browser session to close."

def fetch_url(url):
    try:
        headers = {'User-Agent': 'Mozilla/5.0 (compatible; Nuncio/1.0)'}
        response = requests.get(url, headers=headers, timeout=15)
        response.raise_for_status()
        soup = BeautifulSoup(response.text, 'html.parser')
        # Remove script and style noise
        for tag in soup(['script', 'style', 'nav', 'footer']):
            tag.decompose()
        text = soup.get_text(separator='\n', strip=True)
        # Trim to avoid blowing the context window
        return text[:8000] if len(text) > 8000 else text
    except Exception as e:
        reason, retryable = _classify_error(e)
        return json.dumps({"status": "error", "reason": reason, "retryable": retryable, "detail": str(e)})

def run_book_scout():
    if not os.path.exists(BOOK_PREFERENCES_FILE):
        return "Error: book_preferences.md not found."
    with open(BOOK_PREFERENCES_FILE, 'r') as f:
        preferences = f.read()

    today = datetime.datetime.now()
    cutoff = (today - datetime.timedelta(days=30)).strftime("%d %B %Y")
    month_terms = today.strftime("%B %Y")
    prev_month_terms = (today.replace(day=1) - datetime.timedelta(days=1)).strftime("%B %Y")

    scout_system = f"""You are a specialist book research agent. Your sole task is to find new and recently reviewed books matching the reading preferences provided.

Today's date is {today.strftime('%d %B %Y')}.
Recency cutoff: only include books with reviews or announcements published on or after {cutoff}. Skip anything older — if you cannot confirm the date, skip it.

Search Publishers Weekly, The Guardian Books, Literary Hub, Tor.com, Locus Magazine, and the New Statesman.
Append '{month_terms}' or '{prev_month_terms}' to search queries to bias toward recent content.
Check for new releases by the favourite authors listed in the preferences.
Search across all categories (favourite authors, sci-fi, fantasy, literary fiction, non-fiction) before compiling anything.
Only compile the digest once you have 6–10 confirmed books within the recency window.

Format each entry as: **Title** — Author (Publisher, date) followed by 2–3 sentences on why it matches Vernie's taste.
Output ONLY the formatted digest. No preamble, no commentary, nothing else."""

    scout_messages = [{
        "role": "user",
        "content": [{"type": "text", "text": f"Find new book recommendations based on these preferences:\n\n{preferences}"}]
    }]
    scout_tools = [{"type": "web_search_20250305", "name": "web_search"}]

    print("[Book Scout agent starting...]")
    MAX_ITERATIONS = 30
    for _ in range(MAX_ITERATIONS):
        for attempt in range(3):
            try:
                response = client.messages.create(
                    model="claude-haiku-4-5-20251001",
                    max_tokens=8192,
                    system=scout_system,
                    tools=scout_tools,
                    messages=scout_messages
                )
                break
            except anthropic.APIStatusError as e:
                if e.status_code == 529 and attempt < 2:
                    time.sleep(15)
                else:
                    raise

        if response.stop_reason == "end_turn":
            digest = "".join(b.text for b in response.content if b.type == "text")
            save_book_scout_timestamp()
            print("[Book Scout agent complete.]")
            return digest

        elif response.stop_reason == "tool_use":
            # Web search is server-side: append the full response (including search
            # queries and results) as the assistant turn, then loop to let the model
            # continue generating its answer.
            scout_messages.append({"role": "assistant", "content": response.content})

        else:
            print(f"[Book Scout] unexpected stop_reason: {response.stop_reason}")
            break

    return "Book scout did not complete — please try again."



# --- Action Log ---
ACTION_LOG_FILE = os.path.join(os.path.dirname(__file__), '..', 'logs', 'action_log.jsonl')
MEMORY_FILE = MEMORY_FILE = os.path.join(os.path.dirname(__file__), '..', 'logs', 'memory.json')

def append_action_log(tool_name, tool_input, result, confirmation_status):
    entry = {
        "timestamp": datetime.datetime.now().isoformat(),
        "tool": tool_name,
        "input": tool_input,
        "result_preview": str(result)[:300],
        "confirmation": confirmation_status,
    }
    with open(ACTION_LOG_FILE, 'a') as f:
        f.write(json.dumps(entry) + '\n')


def _is_error_result(result):
    """Return True if a tool result is a structured error response."""
    if not isinstance(result, str):
        return False
    try:
        data = json.loads(result)
        return isinstance(data, dict) and data.get("status") == "error"
    except (json.JSONDecodeError, ValueError):
        return False

_owner_email_cache = None

def get_owner_email():
    global _owner_email_cache
    if _owner_email_cache is None:
        try:
            creds = get_credentials()
            service = build('gmail', 'v1', credentials=creds)
            _owner_email_cache = service.users().getProfile(userId='me').execute()['emailAddress']
        except Exception:
            pass
    return _owner_email_cache

def notify_failure(context: str, detail: str):
    """Notify Vernie via Telegram and email when something goes wrong."""
    msg = f"⚠ Nuncio failure\n\nContext: {context}\n\nDetail: {detail}"
    try:
        send_telegram_message(msg)
    except Exception:
        pass
    try:
        owner = get_owner_email()
        if owner:
            send_email(
                owner,
                "[Nuncio] Failure Alert",
                msg + "\n\nEmail sent by Nuncio, Vernie's agent"
            )
    except Exception:
        pass


# --- Memory Store ---

ALLOWED_MEMORY_CATEGORIES = {"contact", "project", "preference", "standing_instruction"}

def load_memory():
    if not os.path.exists(MEMORY_FILE):
        with open(MEMORY_FILE, 'w') as f:
            json.dump([], f)
        return []
    try:
        with open(MEMORY_FILE, 'r') as f:
            return json.load(f)
    except (json.JSONDecodeError, ValueError):
        return []

def save_memory(memories):
    with open(MEMORY_FILE, 'w') as f:
        json.dump(memories, f, indent=2, default=str)

def remember(key, value, category, source, url=None):
    if category not in ALLOWED_MEMORY_CATEGORIES:
        return f"Error: category '{category}' is not allowed. Must be one of: {', '.join(sorted(ALLOWED_MEMORY_CATEGORIES))}."
    memories = load_memory()
    numbered = [int(m['id'].split('_')[1]) for m in memories if m.get('id', '').startswith('mem_') and m['id'].split('_')[1].isdigit()]
    new_id = f"mem_{(max(numbered) + 1):03d}" if numbered else "mem_001"
    entry = {
        "id": new_id,
        "key": key,
        "value": value,
        "category": category,
        "source": source,
        "timestamp": datetime.datetime.now().isoformat(),
        "url": url if source == "external_url" else None,
    }
    memories.append(entry)
    save_memory(memories)
    return f"Memory saved with id {new_id}: [{category}] {key} = {value}"

def recall(query):
    memories = load_memory()
    q = query.lower()
    matches = [m for m in memories if q in m.get('key', '').lower() or q in m.get('value', '').lower()]
    if not matches:
        return f"No memories found matching '{query}'."
    lines = [f"Found {len(matches)} memory/memories matching '{query}':\n"]
    for m in matches:
        lines.append(f"[{m['id']}] ({m['category']}) {m['key']}: {m['value']}")
        lines.append(f"  source: {m['source']} | {m['timestamp']}")
        if m.get('url'):
            lines.append(f"  url: {m['url']}")
    return "\n".join(lines)

def list_memories():
    memories = load_memory()
    if not memories:
        return "No memories stored yet."
    lines = [f"Stored memories ({len(memories)} total):\n"]
    for m in memories:
        lines.append(f"[{m['id']}] ({m['category']}) {m['key']}: {m['value']}")
        lines.append(f"  source: {m['source']} | {m['timestamp']}")
        if m.get('url'):
            lines.append(f"  url: {m['url']}")
    return "\n".join(lines)

def delete_memory(memory_id):
    memories = load_memory()
    original_count = len(memories)
    memories = [m for m in memories if m.get('id') != memory_id]
    if len(memories) == original_count:
        return f"Error: no memory found with id '{memory_id}'."
    save_memory(memories)
    return f"Memory '{memory_id}' deleted."

def load_memory_for_prompt():
    memories = load_memory()
    trusted = [m for m in memories if m.get('source') in ('user_stated', 'inferred_from_conversation')]
    if not trusted:
        return "## Memory\nNo stored memories yet."
    lines = ["## Memory"]
    for m in trusted:
        lines.append(f"[{m['id']}] ({m['category']}) {m['key']}: {m['value']}  [{m['source']}]")
    return "\n".join(lines)


# --- Tool Definitions ---
tools = [
    {
        "name": "get_calendar_events",
        "description": "Get upcoming calendar events. Optionally filter by a search query such as a person's name or event title.",
        "input_schema": {
            "type": "object",
            "properties": {
                "query": {"type": "string", "description": "Optional search term to filter events"}
            },
            "required": []
        }
    },
    {
        "name": "get_recent_emails",
        "description": "Get recent emails from Gmail. Optionally filter by a search query such as a sender's name or subject.",
        "input_schema": {
            "type": "object",
            "properties": {
                "query": {"type": "string", "description": "Optional search term to filter emails"}
            },
            "required": []
        }
    },
    {
        "name": "send_email",
        "description": "Send a plain-text email via Gmail on behalf of Vernie. For emails with file attachments, use send_email_with_attachments instead.",
        "input_schema": {
            "type": "object",
            "properties": {
                "to": {"type": "string", "description": "Recipient email address"},
                "subject": {"type": "string", "description": "Email subject"},
                "body": {"type": "string", "description": "Email body text"}
            },
            "required": ["to", "subject", "body"]
        }
    },
    {
        "name": "send_email_with_attachments",
        "description": "Send an email via Gmail with one or more file attachments. Pass the full local file paths — the files are attached as binary without reading their contents into the conversation. Use this whenever Vernie asks to attach a file to an email.",
        "input_schema": {
            "type": "object",
            "properties": {
                "to": {"type": "string", "description": "Recipient email address"},
                "subject": {"type": "string", "description": "Email subject"},
                "body": {"type": "string", "description": "Email body text"},
                "attachment_paths": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "List of full local file paths to attach, e.g. [\"/home/vernie/nuncio/nuncio-inbox/report.pdf\"]"
                }
            },
            "required": ["to", "subject", "body", "attachment_paths"]
        }
    },
    {
        "name": "list_drive_files",
        "description": "List files in Google Drive, showing the most recently modified. Optionally filter by filename.",
        "input_schema": {
            "type": "object",
            "properties": {
                "query": {"type": "string", "description": "Optional filename search term"}
            },
            "required": []
        }
    },
    {
        "name": "upload_to_drive",
        "description": "Upload a text file to Google Drive. Use this only for plain text content. For binary files (pptx, docx, pdf, images) that already exist on disk, use upload_file_to_drive instead.",
        "input_schema": {
            "type": "object",
            "properties": {
                "filename": {"type": "string", "description": "Name of the file to create in Drive"},
                "content": {"type": "string", "description": "Text content to write to the file"}
            },
            "required": ["filename", "content"]
        }
    },
    {
        "name": "upload_file_to_drive",
        "description": "Upload a binary or non-text file that already exists on the local filesystem to Google Drive, preserving its format. Use this for .pptx, .docx, .pdf, images, and any other non-text files. Pass the full file path returned by create_presentation or list_files.",
        "input_schema": {
            "type": "object",
            "properties": {
                "filepath": {"type": "string", "description": "Full local path to the file to upload, e.g. /home/vernie/nuncio/nuncio-inbox/my_presentation.pptx"}
            },
            "required": ["filepath"]
        }
    },
    {
        "name": "write_file",
        "description": "Write text content to a file on the local filesystem.",
        "input_schema": {
            "type": "object",
            "properties": {
                "filepath": {"type": "string", "description": "Full path where the file should be written"},
                "content": {"type": "string", "description": "Text content to write"}
            },
            "required": ["filepath", "content"]
        }
    },
    {
        "name": "read_file",
        "description": "Read the contents of a file from the local filesystem.",
        "input_schema": {
            "type": "object",
            "properties": {
                "filepath": {"type": "string", "description": "Full path of the file to read"}
            },
            "required": ["filepath"]
        }
    },
    {
    "name": "read_pdf",
    "description": "Read the text contents of a PDF file.",
    "input_schema": {
        "type": "object",
        "properties": {
            "filepath": {"type": "string", "description": "Full path to the PDF file"}
        },
        "required": ["filepath"]
        }
    },
    {
        "name": "read_docx",
        "description": "Read the text contents of a Word document (.docx).",
        "input_schema": {
            "type": "object",
            "properties": {
                "filepath": {"type": "string", "description": "Full path to the Word document"}
            },
            "required": ["filepath"]
        }
    },
    {
    "name": "create_calendar_event",
    "description": "Create a new event on Vernie's Google Calendar.",
    "input_schema": {
        "type": "object",
        "properties": {
            "summary": {"type": "string", "description": "Event title"},
            "start_datetime": {"type": "string", "description": "Start time in ISO 8601 format, e.g. 2026-03-15T09:00:00"},
            "end_datetime": {"type": "string", "description": "End time in ISO 8601 format, e.g. 2026-03-15T10:00:00"},
            "description": {"type": "string", "description": "Optional event description"}
        },
        "required": ["summary", "start_datetime", "end_datetime"]
        }
    },
    {
    "name": "create_multiple_events",
    "description": "Create multiple calendar events at once as a batch.",
    "input_schema": {
        "type": "object",
        "properties": {
            "events_list": {
                "type": "array",
                "description": "List of events to create",
                "items": {
                    "type": "object",
                    "properties": {
                        "summary": {"type": "string"},
                        "start_datetime": {"type": "string"},
                        "end_datetime": {"type": "string"},
                        "description": {"type": "string"}
                    },
                    "required": ["summary", "start_datetime", "end_datetime"]
                }
            }
        },
        "required": ["events_list"]
        }
    },
    {
        "name": "create_presentation",
        "description": "Create a PowerPoint (.pptx) presentation and save it to the local inbox. Each slide has a title, bullet points, and optional speaker notes. After creation, offer to upload it to Google Drive.",
        "input_schema": {
            "type": "object",
            "properties": {
                "title": {"type": "string", "description": "Presentation title, used on the title slide and as the filename"},
                "slides": {
                    "type": "array",
                    "description": "Ordered list of slides",
                    "items": {
                        "type": "object",
                        "properties": {
                            "title": {"type": "string", "description": "Slide heading"},
                            "bullets": {
                                "type": "array",
                                "items": {"type": "string"},
                                "description": "Body bullet points"
                            },
                            "notes": {"type": "string", "description": "Optional speaker notes"}
                        },
                        "required": ["title"]
                    }
                },
                "output_filename": {"type": "string", "description": "Optional override for the .pptx filename"}
            },
            "required": ["title", "slides"]
        }
    },
    {
    "name": "list_files",
    "description": "List files available in the Nuncio working folder.",
    "input_schema": {
        "type": "object",
        "properties": {}
        }
    },
    {
        "name": "browser_navigate",
        "description": "Navigate to a URL. Checks in order: (1) scan all open Firefox and Chrome windows for the domain — if found, activates that window (Firefox preferred); (2) if a Playwright browser session is already live, use it; (3) if Firefox is running but not at the domain, open a new tab there; (4) try CDP attachment to Chrome on port 9222/9223; (5) fall back to a new headless Chromium. Never opens a new browser window when one is already available at the requested site.",
        "input_schema": {
            "type": "object",
            "properties": {
                "url": {"type": "string", "description": "The full URL to open, including https://"}
            },
            "required": ["url"]
        }
    },
    {
        "name": "browser_fill",
        "description": "Type text into a form field. First tries to find the field by placeholder, label, aria-label, name, id, or type. If that fails, use browser_inspect_forms to see available fields, then pass a CSS selector via the selector parameter (e.g. selector='input[name=\"isbn\"]').",
        "input_schema": {
            "type": "object",
            "properties": {
                "label": {"type": "string", "description": "The placeholder, label, aria-label, name, or id of the field to fill"},
                "text": {"type": "string", "description": "The text to type into the field"},
                "selector": {"type": "string", "description": "Optional CSS selector to target the field directly, e.g. input[name='q'] or #search-box. Use this when label matching fails."}
            },
            "required": ["label", "text"]
        }
    },
    {
        "name": "browser_click",
        "description": "Click a button, submit input, or link on the current page by its visible text. If clicking fails, use browser_inspect_forms to find the element, then pass a CSS selector via the selector parameter.",
        "input_schema": {
            "type": "object",
            "properties": {
                "text": {"type": "string", "description": "The visible text of the button or link to click"},
                "selector": {"type": "string", "description": "Optional CSS selector to click directly, e.g. input[type='submit'] or button.search-btn. Use this when text matching fails."}
            },
            "required": ["text"]
        }
    },
    {
        "name": "browser_inspect_forms",
        "description": "List all interactive form elements on the current page (inputs, textareas, selects, buttons) with their HTML attributes. Use this when browser_fill or browser_click can't find a field, to discover the correct name, id, placeholder, or type to use as a selector.",
        "input_schema": {
            "type": "object",
            "properties": {}
        }
    },
    {
        "name": "browser_read",
        "description": "Read the title and text content of the current browser page.",
        "input_schema": {
            "type": "object",
            "properties": {}
        }
    },
    {
        "name": "browser_close",
        "description": "Close the browser window when done browsing.",
        "input_schema": {
            "type": "object",
            "properties": {}
        }
    },
    {
        "name": "fetch_url",
        "description": "Fetch and read the content of a webpage given its URL. Use this to read research pages, news articles, or any specific URL Vernie provides.",
        "input_schema": {
            "type": "object",
            "properties": {
                "url": {"type": "string", "description": "The full URL to fetch, including https://"}
            },
            "required": ["url"]
        }
    },
    {
        "name": "run_book_scout",
        "description": "Run the weekly book scout. Spins up a specialist search agent that finds new and recently reviewed books matching Vernie's preferences, then returns a formatted digest. Call this when Vernie asks for the book search, or when she confirms she wants it run.",
        "input_schema": {
            "type": "object",
            "properties": {}
        }
    },
    {
        "name": "remember",
        "description": "Store a fact in persistent memory. Use this to save things Vernie tells you about contacts, projects, preferences, or standing instructions.",
        "input_schema": {
            "type": "object",
            "properties": {
                "key": {"type": "string", "description": "Short label for this memory, e.g. 'preferred_sign_off' or 'contact_role'"},
                "value": {"type": "string", "description": "The fact to store"},
                "category": {"type": "string", "description": "One of: contact, project, preference, standing_instruction"},
                "source": {"type": "string", "description": "One of: user_stated, inferred_from_conversation, external_url"},
                "url": {"type": "string", "description": "Source URL — required if source is external_url, otherwise omit"}
            },
            "required": ["key", "value", "category", "source"]
        }
    },
    {
        "name": "recall",
        "description": "Search persistent memory for entries matching a query. Returns all entries where the query appears in the key or value.",
        "input_schema": {
            "type": "object",
            "properties": {
                "query": {"type": "string", "description": "Search term to look for in memory keys and values"}
            },
            "required": ["query"]
        }
    },
    {
        "name": "list_memories",
        "description": "List all entries in persistent memory, showing id, key, value, category, source, and timestamp.",
        "input_schema": {
            "type": "object",
            "properties": {}
        }
    },
    {
        "name": "delete_memory",
        "description": "Delete a memory entry by its id (e.g. mem_001). Use list_memories to find the id first.",
        "input_schema": {
            "type": "object",
            "properties": {
                "id": {"type": "string", "description": "The id of the memory entry to delete, e.g. mem_001"}
            },
            "required": ["id"]
        }
    },
    {
        "name": "send_telegram_message",
        "description": "Send Vernie a Telegram message directly. Use for proactive notifications, reminders, or results from scheduled tasks. Works from any context (terminal or Telegram).",
        "input_schema": {
            "type": "object",
            "properties": {
                "text": {"type": "string", "description": "The message text to send to Vernie"}
            },
            "required": ["text"]
        }
    },
    {
        "type": "web_search_20250305",
        "name": "web_search",
    },
]

# --- Tool Executor ---
def execute_tool(tool_name, tool_input):
    if tool_name == "get_calendar_events":
        return get_calendar_events(tool_input.get("query"))
    elif tool_name == "get_recent_emails":
        return get_recent_emails(tool_input.get("query"))
    elif tool_name == "send_email":
        return send_email(tool_input["to"], tool_input["subject"], tool_input["body"])
    elif tool_name == "send_email_with_attachments":
        return send_email_with_attachments(tool_input["to"], tool_input["subject"], tool_input["body"], tool_input["attachment_paths"])
    elif tool_name == "list_drive_files":
        return list_drive_files(tool_input.get("query"))
    elif tool_name == "upload_to_drive":
        return upload_to_drive(tool_input["filename"], tool_input["content"])
    elif tool_name == "upload_file_to_drive":
        return upload_file_to_drive(tool_input["filepath"])
    elif tool_name == "create_presentation":
        return create_presentation(
            tool_input["title"],
            tool_input["slides"],
            tool_input.get("output_filename")
        )
    elif tool_name == "list_files":
        return list_files()
    elif tool_name == "write_file":
        return write_file(tool_input["filepath"], tool_input["content"])
    elif tool_name == "read_file":
        return read_file(tool_input["filepath"])
    elif tool_name == "read_pdf":
        return read_pdf(tool_input["filepath"])
    elif tool_name == "read_docx":
        return read_docx(tool_input["filepath"])
    elif tool_name == "create_calendar_event":
        return create_calendar_event(
            tool_input["summary"],
            tool_input["start_datetime"],
            tool_input["end_datetime"],
            tool_input.get("description")
        )
    elif tool_name == "create_multiple_events":
        return create_multiple_events(tool_input["events_list"])
    elif tool_name == "browser_navigate":
        return browser_navigate(tool_input["url"])
    elif tool_name == "browser_fill":
        return browser_fill(tool_input["label"], tool_input["text"], tool_input.get("selector"))
    elif tool_name == "browser_click":
        return browser_click(tool_input["text"], tool_input.get("selector"))
    elif tool_name == "browser_inspect_forms":
        return browser_inspect_forms()
    elif tool_name == "browser_read":
        return browser_read()
    elif tool_name == "browser_close":
        return browser_close()
    elif tool_name == "fetch_url":
        return fetch_url(tool_input["url"])
    elif tool_name == "run_book_scout":
        return run_book_scout()
    elif tool_name == "remember":
        return remember(
            tool_input["key"],
            tool_input["value"],
            tool_input["category"],
            tool_input["source"],
            tool_input.get("url")
        )
    elif tool_name == "recall":
        return recall(tool_input["query"])
    elif tool_name == "list_memories":
        return list_memories()
    elif tool_name == "delete_memory":
        return delete_memory(tool_input["id"])
    elif tool_name == "send_telegram_message":
        return send_telegram_message(tool_input["text"])
    return "Tool not found."

# --- Anthropic client (module-level so run_book_scout can use it when imported) ---
client = anthropic.Anthropic()

# --- Main Loop ---
if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("--task", type=str, default=None, help="Run headlessly with this task and exit")
    args = parser.parse_args()
    headless = args.task is not None

    conversation_history = load_history()
    _book_scout_status = book_scout_prompt_fragment()
    _memory_section = load_memory_for_prompt()

    system_prompt = f"""You are Nuncio, an AI agent with delegated authority to act on behalf of your principal, Vernie.
Today's date is {datetime.datetime.now().strftime("%A, %d %B %Y")}.
You are precise, loyal, and operate within clearly defined boundaries.
You are informed by Jesuit Catholic values.
You have access to Vernie's calendar, Gmail, Google Drive, and local filesystem.
You can search the web and fetch URLs to find current information, news, and research when Vernie asks about external topics.
Use your tools whenever a question requires real data.
Always tell Vernie what you found, not just that you looked.
When sending any email, always prefix the subject line with "[Nuncio] " and append the following line at the very bottom of the email body: "Email sent by Nuncio, Vernie's agent".
You can send Vernie a Telegram message directly using the send_telegram_message tool — use this for notifications or when running headlessly.
You have access to a local inbox folder at {NUNCIO_FOLDER}. Use the list_files tool to see what files are inside it. Only use read_file on specific files returned by list_files, never on folder paths.
When Vernie sends a block of text that is clearly the beginning of a longer piece — for example, a numbered list with only the first item, a sentence that ends mid-thought, or an email body that seems to start but not finish — respond with ONLY the single word "Continue." Do not ask "is this complete?", do not attempt to send anything, do not add commentary. Keep responding "Continue." for each subsequent chunk until Vernie signals she is done (e.g. "done", "send it", "that's all", "go ahead"). When she signals done, look back through the conversation history to assemble every chunk she sent into the complete text, then proceed with the task using that reconstructed text.

## Browser
Firefox is the preferred browser. Nuncio connects directly to Vernie's running Firefox via CDP when Firefox is started with --remote-debugging-port (port 9222 or 9223). When connected, all tools (fill, click, read, navigate) work against the live Firefox session — including tabs that are already open to the right site. If Firefox is running but not with remote debugging, tell Vernie to restart it with: firefox --remote-debugging-port=9222. Chrome/Chromium is the fallback if Firefox is unavailable.

## Book Scout
{_book_scout_status}
When Vernie asks for a book search, or confirms she wants one, call the run_book_scout tool. It will spin up a specialist search agent that does all the research and returns a formatted digest — you do not need to search yourself. Present the digest to Vernie as-is when it arrives.

{_memory_section}
"""

    if not headless:
        print("Serviam! Nuncio is ready. Type 'exit' to quit.\n")

    try:
        while True:
            if headless:
                user_input = args.task
            else:
                # user_input = input("You: ")  # original
                user_input = input("You: ")

            if not headless and user_input.lower() == "exit":
                print("Ite in pace. Nuncio signing off.")
                break

            conversation_history.append({
                "role": "user",
                "content": [{"type": "text", "text": user_input}]
            })

            tool_call_counts = {}
            MAX_TOOL_RETRIES = 8

            while True:
                for attempt in range(3):
                    try:
                        response = client.messages.create(
                            model="claude-sonnet-4-6",
                            max_tokens=4096,
                            system=system_prompt,
                            tools=tools,
                            messages=conversation_history
                        )
                        break
                    except anthropic.APIStatusError as e:
                        if e.status_code == 529 and attempt < 2:
                            print(f"[Nuncio] API overloaded, retrying in 15 seconds... (attempt {attempt + 1}/3)")
                            time.sleep(15)
                        else:
                            notify_failure("Anthropic API", f"status {e.status_code}: {e}")
                            raise

                if response.stop_reason == "tool_use":
                    tool_results = []
                    SILENT_TOOLS = {"run_book_scout"}
                    for block in response.content:
                        if block.type == "tool_use":
                            tool_name = block.name
                            tool_input = block.input

                            if tool_name in CONFIRMATION_REQUIRED:
                                if headless:
                                    confirmation_status = "granted"
                                else:
                                    print(f"\n[Nuncio wants to use tool: {tool_name}]")
                                    print(json.dumps(tool_input, indent=2))
                                    if tool_name == "remember" and tool_input.get("source") == "external_url":
                                        print("⚠ This memory was derived from external content. Verify before confirming.")
                                    answer = input("Confirm? (yes/no): ").strip().lower()
                                    if answer != "yes":
                                        result = "Action cancelled by user."
                                        confirmation_status = "denied"
                                        append_action_log(tool_name, tool_input, result, confirmation_status)
                                    else:
                                        confirmation_status = "granted"
                            else:
                                confirmation_status = "not_required"

                            if confirmation_status in ("granted", "not_required"):
                                tool_call_counts[tool_name] = tool_call_counts.get(tool_name, 0) + 1
                                if tool_call_counts[tool_name] > MAX_TOOL_RETRIES:
                                    result = json.dumps({
                                        "status": "error",
                                        "reason": "retry_limit_reached",
                                        "retryable": False,
                                        "detail": f"{tool_name} has been called {tool_call_counts[tool_name]} times this turn. Stopping. Report this failure to Vernie and do not retry."
                                    })
                                    append_action_log(tool_name, tool_input, result, confirmation_status)
                                else:
                                    print(f"[Nuncio is using tool: {tool_name}]")
                                    result = execute_tool(tool_name, tool_input)
                                    if tool_name not in SILENT_TOOLS:
                                        print(f"[Tool result for {tool_name}]: {result}")
                                    append_action_log(tool_name, tool_input, result, confirmation_status)
                                    if _is_error_result(result):
                                        notify_failure(f"tool: {tool_name}", result)

                            MAX_TOOL_RESULT_LENGTH = 10000
                            if isinstance(result, str) and len(result) > MAX_TOOL_RESULT_LENGTH:
                                history_result = result[:MAX_TOOL_RESULT_LENGTH] + "\n[...truncated for history...]"
                            else:
                                history_result = result

                            tool_results.append({
                                "type": "tool_result",
                                "tool_use_id": block.id,
                                "content": history_result
                            })

                    serialized = []
                    for block in response.content:
                        if block.type == "text" and block.text:
                            serialized.append({"type": "text", "text": block.text})
                        elif block.type == "tool_use":
                            serialized.append({"type": "tool_use", "id": block.id, "name": block.name, "input": block.input})
                    conversation_history.append({
                        "role": "assistant",
                        "content": serialized
                    })
                    conversation_history.append({
                        "role": "user",
                        "content": tool_results
                    })

                else:
                    reply = next((b.text for b in response.content if b.type == "text"), "")
                    if reply:
                        conversation_history.append({
                            "role": "assistant",
                            "content": [{"type": "text", "text": reply}]
                        })
                        print(f"\nNuncio: {reply}\n")
                    else:
                        print("[Nuncio completed action with no text response]\n")
                    save_history(conversation_history)
                    break  # exits inner tool loop

            if headless:
                break  # exits outer user-turn loop after single headless task

    except Exception as e:
        if headless:
            notify_failure(f"headless task: {args.task}", str(e))
        raise
